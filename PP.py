import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO


def create_equation_image(equation, filename):
    plt.figure(figsize=(8, 2))
    plt.text(0.5, 0.5, f"${equation}$", horizontalalignment='center', verticalalignment='center', fontsize=20)
    plt.axis('off')
    buf = BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', pad_inches=0.1)
    buf.seek(0)
    return buf


def add_slide_with_image(prs, title, equation, notes=""):
    slide_layout = prs.slide_layouts[5]  # Using a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Create equation image
    equation_img = create_equation_image(equation, f"{title}.png")

    # Add image to slide
    slide.shapes.add_picture(equation_img, Inches(1), Inches(2), width=Inches(8))

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


# Create a new presentation
prs = Presentation()

# Slide 1: Introduction to GANs
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Generative Adversarial Networks (GANs)"
subtitle.text = ""
slide.notes_slide.notes_text_frame.text = "An introduction to Generative Adversarial Networks."

# Slide 2: Discriminator's Role
add_slide_with_image(prs, "Discriminator's Role",
                     "\\max_D V(D, G)",
                     "This slide introduces the objective of the Discriminator (D) in the GAN framework.")

# Slide 3: Real Data Probability
add_slide_with_image(prs, "Real Data Probability",
                     "\\max_D V(D, G) = \\mathbb{E}_{x \\sim p_{\\text{data}}(x)} [\\log D(x)]",
                     "This slide breaks down the part of the Discriminator's objective involving real data probability.")

# Slide 4: Generated Data Probability
add_slide_with_image(prs, "Generated Data Probability",
                     "\\max_D V(D, G) = \\mathbb{E}_{x \\sim p_{\\text{data}}(x)} [\\log D(x)] + \\mathbb{E}_{z \\sim p_z(z)} [\\log (1 - D(G(z)))]",
                     "This slide breaks down the part of the Discriminator's objective involving generated data probability.")

# Slide 5: Generator's Role
add_slide_with_image(prs, "Generator's Role",
                     "\\min_G V(D, G)",
                     "This slide introduces the objective of the Generator (G) in the GAN framework.")

# Slide 6: Generator's Goal
add_slide_with_image(prs, "Generator's Goal",
                     "\\min_G V(D, G) = \\mathbb{E}_{z \\sim p_z(z)} [\\log (1 - D(G(z)))]",
                     "This slide breaks down the Generator's objective involving the probability of generated data being classified as fake.")

# Slide 7: Alternative Generator Objective
add_slide_with_image(prs, "Alternative Generator Objective",
                     "\\max_G \\mathbb{E}_{z \\sim p_z(z)} [\\log D(G(z))]",
                     "This slide presents an alternative objective for the Generator (G).")

# Slide 8: Combined Min-Max Objective
add_slide_with_image(prs, "Combined Min-Max Objective",
                     "\\min_G \\max_D V(D, G)",
                     "This slide introduces the combined min-max objective in the GAN framework.")

# Slide 9: Full Objective Unification
add_slide_with_image(prs, "Full Objective Unification",
                     "\\min_G \\max_D V(D, G) = \\mathbb{E}_{x \\sim p_{\\text{data}}(x)} [\\log D(x)] + \\mathbb{E}_{z \\sim p_z(z)} [\\log (1 - D(G(z)))]",
                     "This slide presents the unified full objective of GANs.")

# Slide 10: Intuitive Explanation
add_slide_with_image(prs, "Intuitive Explanation",
                     "Generator (G): Creates fake data to fool the discriminator.\n"
                     "Discriminator (D): Distinguishes between real and fake data.\n"
                     "Objective: Train G to produce realistic data that D cannot distinguish from real data.",
                     "This slide provides an intuitive explanation of the GAN framework.")

# Save the presentation
prs.save('GAN_Objective_Presentation_with_Equations.pptx')
